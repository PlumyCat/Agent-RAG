import asyncio
import os
import aiofiles
from typing import Dict, Any, List
import logging
from config import Config

logger = logging.getLogger(__name__)

class ContentExtractor:
    """Extract plain text content from various file formats with large file support"""
    
    @staticmethod
    async def get_file_info(file_path: str) -> Dict[str, Any]:
        """Get file information and determine processing strategy"""
        try:
            file_stat = os.stat(file_path)
            file_size_mb = file_stat.st_size / (1024 * 1024)
            
            return {
                "file_path": file_path,
                "file_size_bytes": file_stat.st_size,
                "file_size_mb": round(file_size_mb, 2),
                "is_large_file": file_size_mb > Config.MEMORY_THRESHOLD_MB,
                "processing_strategy": "chunked" if file_size_mb > Config.MEMORY_THRESHOLD_MB else "standard"
            }
        except Exception as e:
            logger.error(f"Failed to get file info for {file_path}: {e}")
            return {"error": str(e)}
    
    # @staticmethod
    # async def extract_pdf_content(file_path: str) -> str:
    #     """Extract all text content from PDF with large file support"""
    #     try:
    #         import PyPDF2
            
    #         file_info = await ContentExtractor.get_file_info(file_path)
    #         is_large = file_info.get("is_large_file", False)
            
    #         content = ""
            
    #         if is_large:
    #             logger.info(f"Processing large PDF file: {file_path}")
    #             # Process in chunks for large files
    #             async with aiofiles.open(file_path, 'rb') as file:
    #                 file_content = await file.read()
                    
    #             # Process PDF from memory
    #             import io
    #             pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                
    #             # Process pages in batches
    #             batch_size = 10
    #             total_pages = len(pdf_reader.pages)
                
    #             for batch_start in range(0, total_pages, batch_size):
    #                 batch_end = min(batch_start + batch_size, total_pages)
                    
    #                 for page_num in range(batch_start, batch_end):
    #                     page = pdf_reader.pages[page_num]
    #                     content += page.extract_text() + "\n"
                    
    #                 # Yield control to prevent blocking
    #                 await asyncio.sleep(0.01)
    #                 logger.info(f"Processed pages {batch_start+1}-{batch_end} of {total_pages}")
    #         else:
    #             # Standard processing for smaller files
    #             with open(file_path, 'rb') as file:
    #                 pdf_reader = PyPDF2.PdfReader(file)
    #                 for page in pdf_reader.pages:
    #                     content += page.extract_text() + "\n"
            
    #         logger.info(f"Extracted {len(content)} characters from PDF: {file_path}")
    #         return content
            
    #     except Exception as e:
    #         logger.error(f"Error extracting PDF content from {file_path}: {e}")
    #         return f"Error extracting PDF content: {str(e)}"

    @staticmethod
    async def extract_pdf_content(file_path: str) -> str:
        """Extract all text content from PDF with enhanced error handling"""
        try:
            import PyPDF2
            
            file_info = await ContentExtractor.get_file_info(file_path)
            is_large = file_info.get("is_large_file", False)
            
            content = ""
            
            if is_large:
                logger.info(f"Processing large PDF file: {file_path}")
                # Process in chunks for large files
                async with aiofiles.open(file_path, 'rb') as file:
                    file_content = await file.read()
                    
                # Process PDF from memory
                import io
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                total_pages = len(pdf_reader.pages)
                
                # Process pages in batches
                batch_size = 10
                
                for batch_start in range(0, total_pages, batch_size):
                    batch_end = min(batch_start + batch_size, total_pages)
                    
                    batch_content = ""  # Accumulate batch content
                    for page_num in range(batch_start, batch_end):
                        try:
                            page = pdf_reader.pages[page_num]
                            page_text = page.extract_text()
                            
                            # Debug: Check if we're getting text
                            if page_text and page_text.strip():
                                batch_content += page_text + "\n"
                                logger.debug(f"Page {page_num + 1}: extracted {len(page_text)} characters")
                            else:
                                logger.warning(f"Page {page_num + 1}: no text extracted")
                                
                        except Exception as e:
                            logger.warning(f"Failed to extract page {page_num + 1}: {e}")
                            continue
                    
                    content += batch_content  # Add batch to total content
                    
                    # Yield control to prevent blocking
                    await asyncio.sleep(0.01)
                    logger.info(f"Processed pages {batch_start+1}-{batch_end} of {total_pages} (batch: {len(batch_content)} chars, total: {len(content)} chars)")
            
            else:
                # Standard processing for smaller files
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text and page_text.strip():
                                content += page_text + "\n"
                        except Exception as e:
                            logger.warning(f"Failed to extract page {page_num + 1}: {e}")
                            continue
            
            logger.info(f"PDF extraction complete: {len(content)} characters from {file_path}")
            
            # If we got very little content, try alternative extraction
            if len(content) < 1000:  # Less than 1000 chars is suspicious
                logger.warning(f"Low content extracted ({len(content)} chars), trying alternative method...")
                try:
                    alternative_content = await ContentExtractor._extract_pdf_alternative(file_path)
                    if len(alternative_content) > len(content):
                        logger.info(f"Alternative method extracted {len(alternative_content)} characters")
                        return alternative_content
                except Exception as e:
                    logger.warning(f"Alternative extraction failed: {e}")
            
            return content
            
        except Exception as e:
            logger.error(f"Error extracting PDF content from {file_path}: {e}")
            return f"Error extracting PDF content: {str(e)}"

    @staticmethod
    async def _extract_pdf_alternative(file_path: str) -> str:
        """Alternative PDF extraction using different libraries"""
        try:
            # Try with pdfplumber (better for complex PDFs)
            try:
                import pdfplumber
                content = ""
                
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text:
                                content += page_text + "\n"
                        except Exception as e:
                            logger.warning(f"pdfplumber failed on page {page_num + 1}: {e}")
                            continue
                
                logger.info(f"pdfplumber extracted {len(content)} characters")
                return content
                
            except ImportError:
                logger.warning("pdfplumber not available, trying pymupdf...")
                
            # Try with pymupdf (fitz)
            try:
                import fitz  # pymupdf
                content = ""
                
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    try:
                        page = doc.load_page(page_num)
                        page_text = page.get_text()
                        if page_text:
                            content += page_text + "\n"
                    except Exception as e:
                        logger.warning(f"pymupdf failed on page {page_num + 1}: {e}")
                        continue
                
                doc.close()
                logger.info(f"pymupdf extracted {len(content)} characters")
                return content
                
            except ImportError:
                logger.warning("pymupdf not available")
                
            return ""
            
        except Exception as e:
            logger.error(f"Alternative PDF extraction failed: {e}")
            return ""
    
    @staticmethod
    async def extract_pdf_content_with_ocr(file_path: str) -> str:
        """Extract content from PDF using OCR for scanned documents"""
        try:
            # First try normal extraction
            content = await ContentExtractor.extract_pdf_content(file_path)
            
            # If we got substantial content, return it
            if len(content.strip()) > 100:
                logger.info(f"PDF contains extractable text: {len(content)} characters")
                return content
            
            # Otherwise, try OCR
            logger.info(f"PDF appears to be scanned, attempting OCR extraction...")
            
            try:
                import pytesseract
                from pdf2image import convert_from_path
                from PIL import Image
                import tempfile
                import os
            except ImportError as e:
                return f"OCR libraries not installed: {e}. Install with: pip install pytesseract pdf2image pillow"
            
            # Convert PDF to images
            logger.info("Converting PDF pages to images...")
            try:
                # Use higher DPI for better OCR accuracy
                pages = convert_from_path(file_path, dpi=300, first_page=1, last_page=None)
                logger.info(f"Converted {len(pages)} pages to images")
            except Exception as e:
                return f"Failed to convert PDF to images: {e}"
            
            # Extract text from each page using OCR
            ocr_content = ""
            total_pages = len(pages)
            
            for page_num, page_image in enumerate(pages, 1):
                try:
                    logger.info(f"Processing page {page_num}/{total_pages} with OCR...")
                    
                    # Convert PIL image to format suitable for OCR
                    page_text = pytesseract.image_to_string(page_image, lang='eng')
                    
                    if page_text.strip():
                        ocr_content += f"\n--- Page {page_num} ---\n"
                        ocr_content += page_text.strip() + "\n"
                        logger.info(f"Page {page_num}: extracted {len(page_text)} characters")
                    else:
                        logger.warning(f"Page {page_num}: no text extracted")
                    
                    # Yield control for async processing
                    if page_num % 5 == 0:
                        await asyncio.sleep(0.1)
                        
                except Exception as e:
                    logger.error(f"OCR failed on page {page_num}: {e}")
                    continue
            
            if ocr_content.strip():
                logger.info(f"OCR extraction complete: {len(ocr_content)} characters")
                return ocr_content
            else:
                return "No text could be extracted from PDF using OCR"
                
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
            return f"OCR extraction failed: {str(e)}"

    @staticmethod
    async def extract_pdf_content_enhanced(file_path: str) -> str:
        """Enhanced PDF extraction that tries multiple methods including OCR"""
        try:
            logger.info(f"Starting enhanced PDF extraction: {file_path}")
            
            # Method 1: Try normal text extraction first
            logger.info("Attempting standard text extraction...")
            standard_content = await ContentExtractor.extract_pdf_content(file_path)
            
            if len(standard_content.strip()) > 100 and not standard_content.startswith("Error"):
                logger.info(f"Standard extraction successful: {len(standard_content)} characters")
                return standard_content
            
            # Method 2: Try OCR extraction
            logger.info("Standard extraction yielded minimal content, trying OCR...")
            ocr_content = await ContentExtractor.extract_pdf_content_with_ocr(file_path)
            
            if len(ocr_content.strip()) > 50 and not ocr_content.startswith("OCR extraction failed"):
                return ocr_content
            
            # If both methods fail
            return f"Unable to extract text from PDF. File may be corrupted, password protected, or contain only images without readable text."
            
        except Exception as e:
            logger.error(f"Enhanced PDF extraction failed: {e}")
            return f"Enhanced PDF extraction failed: {str(e)}"

    @staticmethod
    async def extract_docx_content(file_path: str) -> str:
        """Extract all text content from DOCX with large file support"""
        try:
            from docx import Document
            
            file_info = await ContentExtractor.get_file_info(file_path)
            is_large = file_info.get("is_large_file", False)
            
            if is_large:
                logger.info(f"Processing large DOCX file: {file_path}")
            
            doc = Document(file_path)
            content = ""
            
            # Extract paragraph text
            paragraph_count = len(doc.paragraphs)
            batch_size = 100 if is_large else paragraph_count
            
            for batch_start in range(0, paragraph_count, batch_size):
                batch_end = min(batch_start + batch_size, paragraph_count)
                
                for i in range(batch_start, batch_end):
                    content += doc.paragraphs[i].text + "\n"
                
                if is_large:
                    await asyncio.sleep(0.01)
                    logger.info(f"Processed paragraphs {batch_start+1}-{batch_end} of {paragraph_count}")
            
            # Extract table text
            if doc.tables:
                logger.info(f"Processing {len(doc.tables)} tables")
                for table_idx, table in enumerate(doc.tables):
                    content += f"\n=== Table {table_idx + 1} ===\n"
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        content += " | ".join(row_text) + "\n"
                    
                    if is_large and table_idx % 5 == 0:
                        await asyncio.sleep(0.01)
            
            logger.info(f"Extracted {len(content)} characters from DOCX: {file_path}")
            return content
            
        except Exception as e:
            logger.error(f"Error extracting DOCX content from {file_path}: {e}")
            return f"Error extracting DOCX content: {str(e)}"
    
    @staticmethod
    async def extract_excel_content(file_path: str) -> str:
        """Extract all data from Excel as text with large file support"""
        try:
            import pandas as pd
            
            file_info = await ContentExtractor.get_file_info(file_path)
            is_large = file_info.get("is_large_file", False)
            
            if is_large:
                logger.info(f"Processing large Excel file: {file_path}")
            
            excel_file = pd.ExcelFile(file_path)
            content = ""
            
            for sheet_idx, sheet_name in enumerate(excel_file.sheet_names):
                content += f"\n=== Sheet: {sheet_name} ===\n"
                
                if is_large:
                    # Process large Excel files in chunks
                    chunk_size = 1000
                    chunk_list = []
                    
                    try:
                        for chunk in pd.read_excel(file_path, sheet_name=sheet_name, chunksize=chunk_size):
                            chunk_list.append(chunk)
                            if len(chunk_list) % 5 == 0:
                                await asyncio.sleep(0.01)
                                logger.info(f"Processing chunk {len(chunk_list)} for sheet {sheet_name}")
                        
                        # Combine chunks
                        df = pd.concat(chunk_list, ignore_index=True)
                    except ValueError:
                        # Fallback for sheets that don't support chunking
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                else:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert to string with limited rows for very large sheets
                if len(df) > 10000:
                    content += f"Large sheet with {len(df)} rows, {len(df.columns)} columns\n"
                    content += "First 1000 rows:\n"
                    content += df.head(1000).to_string(max_rows=1000) + "\n"
                    content += f"\n... and {len(df) - 1000} more rows\n"
                else:
                    content += df.to_string() + "\n"
                
                if is_large:
                    await asyncio.sleep(0.01)
                    logger.info(f"Processed sheet {sheet_idx + 1}/{len(excel_file.sheet_names)}: {sheet_name}")
            
            logger.info(f"Extracted {len(content)} characters from Excel: {file_path}")
            return content
            
        except Exception as e:
            logger.error(f"Error extracting Excel content from {file_path}: {e}")
            return f"Error extracting Excel content: {str(e)}"
    
    @staticmethod
    async def extract_csv_content(file_path: str) -> str:
        """Extract all data from CSV as text with large file support"""
        try:
            import pandas as pd
            
            file_info = await ContentExtractor.get_file_info(file_path)
            is_large = file_info.get("is_large_file", False)
            
            if is_large:
                logger.info(f"Processing large CSV file: {file_path}")
                # Process in chunks
                chunk_size = 5000
                chunk_list = []
                
                for chunk in pd.read_csv(file_path, chunksize=chunk_size):
                    chunk_list.append(chunk)
                    if len(chunk_list) % 10 == 0:
                        await asyncio.sleep(0.01)
                        logger.info(f"Processing chunk {len(chunk_list)}")
                
                # Combine chunks
                df = pd.concat(chunk_list, ignore_index=True)
            else:
                df = pd.read_csv(file_path)
            
            # Convert to string with row limits for very large files
            if len(df) > 10000:
                content = f"Large CSV with {len(df)} rows, {len(df.columns)} columns\n"
                content += "Columns: " + ", ".join(df.columns.tolist()) + "\n\n"
                content += "First 1000 rows:\n"
                content += df.head(1000).to_string(max_rows=1000)
                content += f"\n\n... and {len(df) - 1000} more rows"
            else:
                content = df.to_string()
            
            logger.info(f"Extracted {len(content)} characters from CSV: {file_path}")
            return content
            
        except Exception as e:
            logger.error(f"Error extracting CSV content from {file_path}: {e}")
            return f"Error extracting CSV content: {str(e)}"
        
    @staticmethod
    async def extract_pptx_content(file_path: str) -> str:
        """Extract all text content from PowerPoint with large file support"""
        try:
            from pptx import Presentation
            
            file_info = await ContentExtractor.get_file_info(file_path)
            is_large = file_info.get("is_large_file", False)
            
            if is_large:
                logger.info(f"Processing large PowerPoint file: {file_path}")
            
            logger.info(f"Opening PowerPoint presentation: {file_path}")
            prs = Presentation(file_path)
            content = ""
            
            total_slides = len(prs.slides)
            logger.info(f"Found {total_slides} slides in presentation")
            
            # Process slides in batches for large files
            batch_size = 5 if is_large else total_slides
            
            for batch_start in range(0, total_slides, batch_size):
                batch_end = min(batch_start + batch_size, total_slides)
                batch_content = ""
                
                for slide_idx in range(batch_start, batch_end):
                    slide = prs.slides[slide_idx]
                    slide_number = slide_idx + 1
                    
                    try:
                        slide_text = f"\n=== Slide {slide_number} ===\n"
                        
                        # Extract slide title
                        title_text = ""
                        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                            if hasattr(slide.shapes.title, 'text'):
                                title_text = slide.shapes.title.text.strip()
                                if title_text:
                                    slide_text += f"Title: {title_text}\n\n"
                        
                        # Extract content from all shapes
                        shape_texts = []
                        for shape in slide.shapes:
                            shape_content = ""
                            
                            # Handle text frames (most common)
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                if hasattr(shape, 'text') and shape.text.strip():
                                    shape_content = shape.text.strip()
                            
                            # Handle shapes with text attribute directly
                            elif hasattr(shape, 'text') and shape.text.strip():
                                shape_content = shape.text.strip()
                            
                            # Handle tables
                            elif hasattr(shape, 'has_table') and shape.has_table:
                                try:
                                    table_content = "\n[TABLE]\n"
                                    table = shape.table
                                    for row_idx, row in enumerate(table.rows):
                                        row_data = []
                                        for cell in row.cells:
                                            cell_text = cell.text.strip() if hasattr(cell, 'text') else ""
                                            row_data.append(cell_text)
                                        table_content += " | ".join(row_data) + "\n"
                                    table_content += "[/TABLE]\n"
                                    shape_content = table_content
                                except Exception as e:
                                    logger.warning(f"Failed to extract table from slide {slide_number}: {e}")
                            
                            # Add non-empty content
                            if shape_content and shape_content not in shape_texts:
                                # Skip if it's the same as title to avoid duplication
                                if shape_content != title_text:
                                    shape_texts.append(shape_content)
                        
                        # Add all shape content
                        if shape_texts:
                            slide_text += "\n".join(shape_texts) + "\n"
                        
                        # Extract slide notes
                        try:
                            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                                notes_text = ""
                                for shape in slide.notes_slide.shapes:
                                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                        if hasattr(shape, 'text') and shape.text.strip():
                                            notes_text += shape.text.strip() + " "
                                
                                if notes_text.strip():
                                    slide_text += f"\nNotes: {notes_text.strip()}\n"
                        except Exception as e:
                            logger.debug(f"No notes or failed to extract notes from slide {slide_number}: {e}")
                        
                        batch_content += slide_text
                        logger.debug(f"Slide {slide_number}: extracted {len(slide_text)} characters")
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract content from slide {slide_number}: {e}")
                        batch_content += f"\n=== Slide {slide_number} ===\n[Error extracting slide content]\n"
                
                content += batch_content
                
                # Yield control for large files
                if is_large:
                    await asyncio.sleep(0.01)
                    logger.info(f"Processed slides {batch_start+1}-{batch_end} of {total_slides}")
            
            # Extract presentation properties if available
            try:
                if hasattr(prs, 'core_properties'):
                    props = prs.core_properties
                    metadata_text = "\n=== Presentation Metadata ===\n"
                    
                    if hasattr(props, 'title') and props.title:
                        metadata_text += f"Presentation Title: {props.title}\n"
                    if hasattr(props, 'author') and props.author:
                        metadata_text += f"Author: {props.author}\n"
                    if hasattr(props, 'subject') and props.subject:
                        metadata_text += f"Subject: {props.subject}\n"
                    if hasattr(props, 'created') and props.created:
                        metadata_text += f"Created: {props.created}\n"
                    
                    if metadata_text != "\n=== Presentation Metadata ===\n":
                        content = metadata_text + content
            except Exception as e:
                logger.debug(f"Failed to extract metadata: {e}")
            
            logger.info(f"Extracted {len(content)} characters from PowerPoint: {file_path}")
            return content
            
        except ImportError:
            return "Error: python-pptx library not installed. Install with: pip install python-pptx"
        except Exception as e:
            logger.error(f"Error extracting PowerPoint content from {file_path}: {e}")
            return f"Error extracting PowerPoint content: {str(e)}"

    @staticmethod
    async def extract_ppt_content(file_path: str) -> str:
        """Handle legacy PPT files with guidance"""
        try:
            # For old .ppt files, provide helpful error message
            return """Error: Legacy PPT format not supported directly.
            
            Options to convert PPT to PPTX:
            1. Open in PowerPoint and save as .pptx
            2. Use online converters like CloudConvert
            3. Use LibreOffice to convert: libreoffice --headless --convert-to pptx file.ppt

            Note: python-pptx only supports the modern PPTX format (PowerPoint 2007+)."""
                    
        except Exception as e:
                    logger.error(f"Error handling PPT file {file_path}: {e}")
                    return f"Error handling PPT file: {str(e)}"
        
    @staticmethod
    async def preprocess_image_for_ocr(image_path: str, enhancement_level: str = "standard") -> str:
        """Preprocess image for better OCR results with multiple enhancement levels"""
        try:
            import cv2
            import numpy as np
            from PIL import Image, ImageEnhance, ImageFilter
            import tempfile
            import os
            
            logger.info(f"Preprocessing image for OCR: {image_path} (level: {enhancement_level})")
            
            # Read image with OpenCV
            img = cv2.imread(image_path)
            if img is None:
                # Fallback to PIL if OpenCV fails
                pil_img = Image.open(image_path)
                img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
            
            original_img = img.copy()
            
            # Apply preprocessing based on enhancement level
            if enhancement_level == "light":
                processed_img = await ContentExtractor._light_preprocessing(img)
            elif enhancement_level == "standard":
                processed_img = await ContentExtractor._standard_preprocessing(img)
            elif enhancement_level == "aggressive":
                processed_img = await ContentExtractor._aggressive_preprocessing(img)
            else:
                processed_img = img
            
            # Save preprocessed image to temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            temp_path = temp_file.name
            temp_file.close()
            
            cv2.imwrite(temp_path, processed_img)
            logger.info(f"Preprocessed image saved to: {temp_path}")
            
            return temp_path
            
        except ImportError:
            logger.warning("OpenCV not available, using original image")
            return image_path
        except Exception as e:
            logger.error(f"Image preprocessing failed: {e}")
            return image_path

    @staticmethod
    async def _light_preprocessing(img):
        """Light preprocessing - minimal enhancement"""
        import cv2
        
        # Convert to grayscale
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        else:
            gray = img
        
        # Slight denoising
        denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
        
        return denoised

    @staticmethod
    async def _standard_preprocessing(img):
        """Standard preprocessing - balanced enhancement"""
        import cv2
        import numpy as np
        
        # Convert to grayscale
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        else:
            gray = img
        
        # Resize if too small (OCR works better with higher DPI)
        height, width = gray.shape[:2]
        if width < 1000:
            scale_factor = 1000 / width
            new_width = int(width * scale_factor)
            new_height = int(height * scale_factor)
            gray = cv2.resize(gray, (new_width, new_height), interpolation=cv2.INTER_CUBIC)
        
        # Denoising
        denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
        
        # Adaptive thresholding
        thresh = cv2.adaptiveThreshold(
            denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        
        # Morphological operations to clean up
        kernel = np.ones((1, 1), np.uint8)
        cleaned = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
        
        return cleaned

    @staticmethod
    async def _aggressive_preprocessing(img):
        """Aggressive preprocessing - maximum enhancement for difficult images"""
        import cv2
        import numpy as np
        
        # Convert to grayscale
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        else:
            gray = img
        
        # Resize for better OCR (ensure minimum 300 DPI equivalent)
        height, width = gray.shape[:2]
        if width < 1200:
            scale_factor = 1200 / width
            new_width = int(width * scale_factor)
            new_height = int(height * scale_factor)
            gray = cv2.resize(gray, (new_width, new_height), interpolation=cv2.INTER_CUBIC)
        
        # Strong denoising
        denoised = cv2.fastNlMeansDenoising(gray, None, 15, 7, 21)
        
        # Deskew detection and correction
        try:
            coords = np.column_stack(np.where(denoised > 0))
            if len(coords) > 0:
                angle = cv2.minAreaRect(coords)[-1]
                if angle < -45:
                    angle = -(90 + angle)
                else:
                    angle = -angle
                
                # Only apply rotation if angle is significant
                if abs(angle) > 0.5:
                    (h, w) = denoised.shape[:2]
                    center = (w // 2, h // 2)
                    M = cv2.getRotationMatrix2D(center, angle, 1.0)
                    denoised = cv2.warpAffine(denoised, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
        except Exception as e:
            logger.debug(f"Deskew failed, continuing: {e}")
        
        # Contrast enhancement
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(denoised)
        
        # Multiple thresholding approaches and combine
        thresh1 = cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
        thresh2 = cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY, 11, 2)
        
        # Combine thresholds
        combined = cv2.bitwise_and(thresh1, thresh2)
        
        # Morphological operations
        kernel = np.ones((2, 2), np.uint8)
        cleaned = cv2.morphologyEx(combined, cv2.MORPH_CLOSE, kernel)
        cleaned = cv2.morphologyEx(cleaned, cv2.MORPH_OPEN, kernel)
        
        return cleaned

    @staticmethod
    async def extract_image_content(file_path: str) -> str:
        """Extract text from images using OCR with multiple enhancement levels"""
        try:
            import pytesseract
            from PIL import Image
            import os
            
            file_info = await ContentExtractor.get_file_info(file_path)
            file_size_mb = file_info.get("file_size_mb", 0)
            
            logger.info(f"Starting OCR extraction from image: {file_path} ({file_size_mb:.1f}MB)")
            
            # Try multiple enhancement levels and OCR configurations
            enhancement_levels = ["light", "standard", "aggressive"]
            ocr_configs = [
                "--psm 3",  # Default: Fully automatic page segmentation
                "--psm 6",  # Assume uniform block of text
                "--psm 8",  # Treat as single word
                "--psm 11", # Sparse text
                "--psm 13"  # Raw line. Treat as single text line
            ]
            
            best_result = ""
            best_confidence = 0
            
            for enhancement in enhancement_levels:
                try:
                    # Preprocess image
                    processed_image_path = await ContentExtractor.preprocess_image_for_ocr(
                        file_path, enhancement
                    )
                    
                    # Try different OCR configurations
                    for config in ocr_configs:
                        try:
                            # Extract text with current configuration
                            with Image.open(processed_image_path) as img:
                                # Basic extraction
                                text = pytesseract.image_to_string(img, config=config, lang='eng')
                                
                                # Get confidence data
                                try:
                                    data = pytesseract.image_to_data(img, config=config, output_type=pytesseract.Output.DICT)
                                    confidences = [int(x) for x in data['conf'] if int(x) > 0]
                                    avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                                except:
                                    avg_confidence = len(text.strip()) * 10  # Rough confidence based on text length
                                
                                # Keep best result
                                if len(text.strip()) > len(best_result.strip()) or avg_confidence > best_confidence:
                                    best_result = text
                                    best_confidence = avg_confidence
                                    logger.info(f"Better result with {enhancement} + {config}: {avg_confidence:.1f}% confidence, {len(text)} chars")
                        
                        except Exception as e:
                            logger.debug(f"OCR config {config} failed: {e}")
                            continue
                    
                    # Clean up preprocessed file
                    if processed_image_path != file_path:
                        try:
                            os.unlink(processed_image_path)
                        except:
                            pass
                            
                except Exception as e:
                    logger.warning(f"Enhancement level {enhancement} failed: {e}")
                    continue
            
            # Format result
            if best_result.strip():
                content = f"=== Image OCR Extraction ===\n"
                content += f"File: {os.path.basename(file_path)}\n"
                content += f"Size: {file_size_mb:.1f}MB\n"
                content += f"Confidence: {best_confidence:.1f}%\n\n"
                content += "Extracted Text:\n"
                content += best_result.strip()
                
                logger.info(f"OCR extraction complete: {len(content)} characters, {best_confidence:.1f}% confidence")
                return content
            else:
                return f"No text could be extracted from image: {os.path.basename(file_path)}"
                
        except ImportError:
            return "Error: pytesseract not installed. Install with: pip install pytesseract"
        except Exception as e:
            logger.error(f"Image OCR extraction failed from {file_path}: {e}")
            return f"Error extracting text from image: {str(e)}"

    @staticmethod
    async def extract_image_batch(image_paths: List[str]) -> str:
        """Process multiple images in batch for better efficiency"""
        try:
            logger.info(f"Starting batch OCR processing for {len(image_paths)} images")
            
            all_content = "=== Batch Image OCR Results ===\n\n"
            successful_extractions = 0
            
            for i, image_path in enumerate(image_paths, 1):
                try:
                    logger.info(f"Processing image {i}/{len(image_paths)}: {os.path.basename(image_path)}")
                    
                    content = await ContentExtractor.extract_image_content(image_path)
                    
                    if not content.startswith("Error") and not content.startswith("No text"):
                        successful_extractions += 1
                        all_content += f"\n--- Image {i}: {os.path.basename(image_path)} ---\n"
                        all_content += content.replace("=== Image OCR Extraction ===\n", "") + "\n"
                    else:
                        all_content += f"\n--- Image {i}: {os.path.basename(image_path)} ---\n"
                        all_content += f"Failed: {content}\n"
                    
                    # Yield control for async processing
                    if i % 3 == 0:
                        await asyncio.sleep(0.1)
                        
                except Exception as e:
                    logger.error(f"Failed to process image {image_path}: {e}")
                    all_content += f"\n--- Image {i}: {os.path.basename(image_path)} ---\n"
                    all_content += f"Error: {str(e)}\n"
            
            # Add summary
            summary = f"\n=== Batch Summary ===\n"
            summary += f"Total images: {len(image_paths)}\n"
            summary += f"Successful extractions: {successful_extractions}\n"
            summary += f"Success rate: {(successful_extractions/len(image_paths)*100):.1f}%\n"
            
            all_content = summary + all_content
            
            logger.info(f"Batch processing complete: {successful_extractions}/{len(image_paths)} successful")
            return all_content
            
        except Exception as e:
            logger.error(f"Batch image processing failed: {e}")
            return f"Batch processing error: {str(e)}"
    
    @staticmethod 
    async def extract_text_content(file_path: str) -> str:
        """Extract content from plain text files (.txt, .md) with large file support"""
        try:
            file_info = await ContentExtractor.get_file_info(file_path)
            is_large = file_info.get("is_large_file", False)
            
            if is_large:
                logger.info(f"Processing large text file: {file_path}")
            
            # Read file with proper encoding detection
            encodings_to_try = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252']
            content = ""
            
            for encoding in encodings_to_try:
                try:
                    async with aiofiles.open(file_path, 'r', encoding=encoding) as file:
                        if is_large:
                            # Read in chunks for large files
                            chunk_size = 1024 * 1024  # 1MB chunks
                            while True:
                                chunk = await file.read(chunk_size)
                                if not chunk:
                                    break
                                content += chunk
                                await asyncio.sleep(0.01)  # Yield control
                        else:
                            content = await file.read()
                    
                    logger.info(f"Successfully read text file with {encoding} encoding: {len(content)} characters")
                    return content
                    
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    logger.error(f"Error reading file with {encoding}: {e}")
                    continue
                    
            raise Exception(f"Failed to decode file with any of the attempted encodings: {encodings_to_try}")
                    
        except Exception as e:
            logger.error(f"Failed to extract text content from {file_path}: {e}")
            raise Exception(f"Text extraction failed: {str(e)}")
            
    @staticmethod
    async def extract_content(file_path: str) -> tuple[str, str]:
        """Extract content from any supported file type including images"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        extractors = {
            '.pdf': ContentExtractor.extract_pdf_content_enhanced,
            '.docx': ContentExtractor.extract_docx_content,
            '.xlsx': ContentExtractor.extract_excel_content,
            '.xls': ContentExtractor.extract_excel_content,
            '.csv': ContentExtractor.extract_csv_content,
            '.pptx': ContentExtractor.extract_pptx_content,
            '.ppt': ContentExtractor.extract_ppt_content,
            '.pptm': ContentExtractor.extract_pptx_content,
            '.potx': ContentExtractor.extract_pptx_content,
            # Text formats
            '.txt': ContentExtractor.extract_text_content,
            '.md': ContentExtractor.extract_text_content,
            # Image formats
            '.jpg': ContentExtractor.extract_image_content,
            '.jpeg': ContentExtractor.extract_image_content,
            '.png': ContentExtractor.extract_image_content,
            '.bmp': ContentExtractor.extract_image_content,
            '.tiff': ContentExtractor.extract_image_content,
            '.tif': ContentExtractor.extract_image_content,
            '.gif': ContentExtractor.extract_image_content,
            '.webp': ContentExtractor.extract_image_content,
        }
        
        file_type_map = {
            '.pdf': 'PDF',
            '.docx': 'DOCX',
            '.xlsx': 'Excel',
            '.xls': 'Excel',
            '.csv': 'CSV',
            '.pptx': 'PowerPoint',
            '.ppt': 'PowerPoint',
            '.pptm': 'PowerPoint',
            '.potx': 'PowerPoint',
            '.txt': 'Text',
            '.md': 'Markdown',
            '.jpg': 'Image',
            '.jpeg': 'Image',
            '.png': 'Image',
            '.bmp': 'Image',
            '.tiff': 'Image',
            '.tif': 'Image',
            '.gif': 'Image',
            '.webp': 'Image',
        }
        
        if file_ext in extractors:
            content = await extractors[file_ext](file_path)
            file_type = file_type_map[file_ext]
            return content, file_type
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
